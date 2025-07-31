import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.vectorstores import FAISS
from dotenv import load_dotenv
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
import os
import pandas as pd
from pptx import Presentation
import time
from typing import List, Dict, Any
import io

load_dotenv()
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

# Supported file extensions
SUPPORTED_EXTENSIONS = ['pdf', 'docx', 'txt', 'csv', 'xlsx', 'xls', 'pptx', 'ppt', 'r', 'py', 'java', 'c', 'cpp']

def get_java_text(java_file):
    """Read Java source code file"""
    text = java_file.getvalue().decode('utf-8')
    return text

def get_python_text(py_file):
    """Read Python source code file"""
    text = py_file.getvalue().decode('utf-8')
    return text

def get_r_text(r_file):
    """Read R source code file"""
    text = r_file.getvalue().decode('utf-8')
    return text

def get_c_text(c_file):
    """Read C source code file"""
    text = c_file.getvalue().decode('utf-8')
    return text

def get_cpp_text(cpp_file):
    """Read C++ source code file"""
    text = cpp_file.getvalue().decode('utf-8')
    return text

def combine_text(text_list):
    return "\n".join(text_list)

def get_pdf_text(pdf_file):
    text = ""
    pdf_reader = PdfReader(pdf_file)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_word_text(docx_file):
    document = Document(docx_file)
    text = "\n".join([paragraph.text for paragraph in document.paragraphs])
    return text

def read_text_file(txt_file):
    text = txt_file.getvalue().decode('utf-8')
    return text

def get_csv_text(csv_file):
    """Extract text content from CSV files"""
    try:
        # Read CSV file
        df = pd.read_csv(csv_file)
        
        # Convert DataFrame to readable text format
        text_content = []
        
        # Add column headers
        headers = "Column Headers: " + ", ".join(df.columns.tolist())
        text_content.append(headers)
        text_content.append("\n" + "="*50 + "\n")
        
        # Add data summary
        text_content.append(f"Total Rows: {len(df)}")
        text_content.append(f"Total Columns: {len(df.columns)}")
        text_content.append("\n" + "-"*30 + "\n")
        
        # Add sample data (first few rows)
        text_content.append("Sample Data:")
        for idx, row in df.head(10).iterrows():  # Show first 10 rows
            row_text = f"Row {idx + 1}: " + " | ".join([f"{col}: {str(val)}" for col, val in row.items()])
            text_content.append(row_text)
        
        # Add data types info
        text_content.append("\n" + "-"*30 + "\n")
        text_content.append("Data Types:")
        for col, dtype in df.dtypes.items():
            text_content.append(f"{col}: {dtype}")
        
        # Add statistical summary for numeric columns
        numeric_cols = df.select_dtypes(include=['number']).columns
        if len(numeric_cols) > 0:
            text_content.append("\n" + "-"*30 + "\n")
            text_content.append("Numeric Column Statistics:")
            for col in numeric_cols:
                stats = df[col].describe()
                text_content.append(f"{col} - Mean: {stats['mean']:.2f}, Min: {stats['min']}, Max: {stats['max']}")
        
        return "\n".join(text_content)
        
    except Exception as e:
        return f"Error reading CSV file: {str(e)}"

def get_excel_text(excel_file):
    """Extract text content from Excel files"""
    try:
        # Read all sheets from Excel file
        excel_data = pd.read_excel(excel_file, sheet_name=None)
        
        text_content = []
        text_content.append(f"Excel File contains {len(excel_data)} sheet(s)")
        text_content.append("="*50)
        
        for sheet_name, df in excel_data.items():
            text_content.append(f"\nSheet: {sheet_name}")
            text_content.append("-" * 30)
            
            # Add column headers
            headers = "Columns: " + ", ".join(df.columns.tolist())
            text_content.append(headers)
            
            # Add basic info
            text_content.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
            
            # Add sample data (first 5 rows for each sheet)
            if not df.empty:
                text_content.append("\nSample Data:")
                for idx, row in df.head(5).iterrows():
                    row_text = f"Row {idx + 1}: " + " | ".join([f"{col}: {str(val)}" for col, val in row.items()])
                    text_content.append(row_text)
            
            # Add statistical summary for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                text_content.append(f"\nNumeric Columns Summary:")
                for col in numeric_cols:
                    if not df[col].empty:
                        stats = df[col].describe()
                        text_content.append(f"{col} - Mean: {stats['mean']:.2f}, Min: {stats['min']}, Max: {stats['max']}")
        
        return "\n".join(text_content)
        
    except Exception as e:
        return f"Error reading Excel file: {str(e)}"

def get_ppt_text(ppt_file):
    """Extract text content from PowerPoint files"""
    try:
        presentation = Presentation(ppt_file)
        text_content = []
        
        text_content.append(f"PowerPoint Presentation with {len(presentation.slides)} slides")
        text_content.append("="*50)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"\nSlide {slide_num}:")
            text_content.append("-" * 20)
            
            # Extract text from all shapes in the slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_content.extend(slide_text)
            else:
                text_content.append("(No text content found)")
        
        return "\n".join(text_content)
        
    except Exception as e:
        return f"Error reading PowerPoint file: {str(e)}"

def is_supported_file(filename):
    """Check if file extension is supported"""
    file_ext = filename.lower().split('.')[-1]
    return file_ext in SUPPORTED_EXTENSIONS

def get_file_text(file):
    """Route file to appropriate text extraction function based on extension"""
    filename = file.name.lower()
    
    if not is_supported_file(filename):
        return f"Unsupported file format. Please upload only: {', '.join(SUPPORTED_EXTENSIONS)}"
    
    try:
        if filename.endswith('.pdf'):
            return get_pdf_text(file)
        elif filename.endswith('.docx'):
            return get_word_text(file)
        elif filename.endswith('.txt'):
            return read_text_file(file)
        elif filename.endswith('.csv'):
            return get_csv_text(file)
        elif filename.endswith(('.xlsx', '.xls')):
            return get_excel_text(file)
        elif filename.endswith(('.pptx', '.ppt')):
            return get_ppt_text(file)
        elif filename.endswith('.pdf'):
            return get_python_text(file)
        elif filename.endswith('.java'):
            return get_java_text(file)
        elif filename.endswith('.r'):
            return get_r_text(file)
        elif filename.endswith('.c'):
            return get_c_text(file)
        elif filename.endswith('.cpp'):
            return get_cpp_text(file)
        
        
        else:
            return f"Unsupported file format: {filename}"
    except Exception as e:
        return f"Error processing {filename}: {str(e)}"

def get_chunks_with_metadata(text: str, document_id: str, document_name: str, 
                           ) -> List[Dict[str, Any]]:
    """
    Create chunks with document metadata, including code summaries for programming files
    
    Args:
        text: The document text to chunk
        document_id: Unique identifier for the document
        document_name: Name/path of the document    
    Returns:
        List of chunks with metadata
    """
    
    # Define file type categories
    regular_formats = ['pdf', 'docx', 'txt', 'csv', 'xlsx', 'xls', 'pptx', 'ppt']
    code_formats = ['r', 'py', 'java', 'c', 'cpp']
    
    # Extract file extension
    file_extension = document_name.split('.')[-1].lower() if '.' in document_name else document_name.lower()
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = text_splitter.split_text(text)
    
    # Initialize Gemini if dealing with code files
    if file_extension in code_formats:
        # Configure Gemini API
        api_key = os.getenv('GOOGLE_API_KEY')
        if not api_key:
            raise ValueError("Google API key is required for code file processing. "
                           "Provide it as parameter or set GOOGLE_API_KEY environment variable.")
        
        llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash", temperature=0)
    
    # Create metadata for each chunk
    chunks_with_metadata = []
    
    for i, chunk in enumerate(chunks):
        # Base metadata
        metadata = {
            'document_id': document_id,
            'document_name': document_name,
            'chunk_index': i,
            'total_chunks': len(chunks),
            'file_type': file_extension
        }
        
        # Generate summary for code files
        if file_extension in code_formats:
            try:
                summary = generate_code_summary(chunk, file_extension, llm)
                metadata['code_summary'] = summary
                metadata['has_code_summary'] = True
                
                # Add a small delay to respect API rate limits
                time.sleep(5)
                
            except Exception as e:
                print(f"Warning: Failed to generate summary for chunk {i}: {str(e)}")
                metadata['code_summary'] = "Summary generation failed"
                metadata['has_code_summary'] = False
                metadata['summary_error'] = str(e)
        else:
            metadata['has_code_summary'] = False
        
        # For code files, create summary-enhanced text
        if file_extension in code_formats and metadata.get('has_code_summary', False):
            # Option A: Prepend summary for better search
            enhanced_text = f"SUMMARY: {metadata['code_summary']}\n\nCODE:\n{chunk}"
            
            # Option B: Replace with summary-focused format (uncomment to use)
            # enhanced_text = f"{metadata['code_summary']}\n\n{chunk}"
            
        else:
            enhanced_text = chunk

        chunks_with_metadata.append({
            'text': enhanced_text,
            'metadata': metadata
        })
    
    return chunks_with_metadata

def generate_code_summary(code_chunk: str, file_extension: str, llm) -> str:
    """
    Generate a summary description of a code chunk using Google Gemini
    
    Args:
        code_chunk: The code text to summarize
        file_extension: The programming language file extension
        model: Configured Gemini model instance
    
    Returns:
        Summary description of the code chunk
    """
    
    # Language mapping for better prompts
    language_map = {
        'py': 'Python',
        'r': 'R',
        'java': 'Java',
        'c': 'C',
        'cpp': 'C++'
    }
    
    language = language_map.get(file_extension, file_extension.upper())
    
    prompt = f"""
Analyze this {language} code chunk and provide a concise summary (2-3 sentences) that describes:
1. What the code does (main functionality)
2. Key functions, classes, or important variables if present
3. The purpose or role of this code section

Code chunk:
```{file_extension}
{code_chunk}
```

Provide only the summary without any additional formatting or explanations.
"""
    
    try:
        response = llm.invoke(prompt)
        return response.content.strip()
    except Exception as e:
        return f"Error generating summary: {str(e)}"
    
def create_vector_store_with_metadata(all_chunks_with_metadata):
    """Create vector store with document metadata"""
    texts = [chunk['text'] for chunk in all_chunks_with_metadata]
    metadatas = [chunk['metadata'] for chunk in all_chunks_with_metadata]
    
    vector_store = FAISS.from_texts(texts, embedding=embeddings, metadatas=metadatas)
    vector_store.save_local("faiss_db")
    
    return vector_store

def get_top_retrieval_results(retriever, question, k=5):
    """Enhanced version with better error handling and debugging"""
    try:
        # Search for relevant chunks
        docs = retriever.similarity_search_with_score(question, k=k)
        # print(f"Retrieved {len(docs)} documents for question: '{question[:50]}...'")
        
        results = []
        for i, (doc, score) in enumerate(docs):
            try:
                # Convert numpy.float32 to regular float
                score_value = float(score)
                
                # Extract all available metadata
                metadata = doc.metadata if hasattr(doc, 'metadata') else {}
                
                result = {
                    'content': doc.page_content,
                    'score': score_value,
                    'source': metadata.get('document_name', 'Unknown'),
                    'document_id': metadata.get('document_id', 'Unknown'),
                    'chunk_index': metadata.get('chunk_index', 0),
                    'total_chunks': metadata.get('total_chunks', 'Unknown'),
                    'doc_uuid': getattr(doc, 'id', f'doc_{i}'),
                    'rank': i + 1  # Add ranking information
                }
                results.append(result)
                
                # print(f"Document {i+1}: {metadata.get('document_name', 'Unknown')} "
                #       f"(chunk {metadata.get('chunk_index', 0)}/{metadata.get('total_chunks', '?')}) "
                #       f"- Score: {score_value:.4f}")
                
            except Exception as doc_error:
                print(f"Error processing document {i}: {str(doc_error)}")
                continue
        
        print(f"Successfully processed {len(results)} documents")
        return results
        
    except Exception as e:
        error_msg = f"Error searching documents: {str(e)}"
        print(error_msg)
        
        # If using streamlit
        if 'st' in globals():
            st.error(error_msg)
        return []

def get_conversational_chain(retriever, ques, chat_history):
    """
    Get conversational response using retriever and chat history
    
    Args:
        retriever: Document retriever object
        ques: User question string
        chat_history: List of previous chat messages
    
    Returns:
        str: Generated response from the LLM
    """
    try:
        llm = ChatGoogleGenerativeAI(model="gemini-2.0-flash", temperature=0)
        
        # Get relevant documents with scores
        docs_with_scores = retriever.similarity_search_with_score(ques)
        # print(f"Retrieved {len(docs_with_scores)} documents for question: '{ques[:50]}...'")
        
        # Extract context from documents (unpack the tuples correctly)
        context_parts = []
        for doc, score in docs_with_scores:
            # Only include the page content, not metadata info
            context_parts.append(doc.page_content)
        
        context = "\n\n".join(context_parts)
        
        # Format chat history for context
        chat_context = ""
        if chat_history:
            chat_context = "\n".join([
                f"{'User' if msg['role'] == 'user' else 'Assistant'}: {msg['content']}" 
                for msg in chat_history[-6:]  # Last 6 messages for context
            ])
        
        # Create a comprehensive prompt
        prompt_template = """You are a helpful medical information assistant. Answer the question as detailed as possible from the provided context.
Make sure to provide all the details and be accurate. Consider the chat history for context and continuity.
If the answer is not in the provided context or chat history, just say "answer is not available in the context", don't provide the wrong answer.

Chat History:
{chat_history}

Context from documents:
{context}

Question: {question}

Instructions:
- Provide a comprehensive answer based on the context
- Include relevant details from the source documents
- Maintain consistency with previous conversation
- If information is incomplete, state what information is available
- Do not make up information not present in the context

Answer:"""
        
        prompt = prompt_template.format(
            chat_history=chat_context,
            context=context,
            question=ques
        )
        
        # Get response directly from LLM
        response = llm.invoke(prompt)
        
        # print(f"Generated response length: {len(response.content)} characters")
        return response.content
        
    except Exception as e:
        error_msg = f"Error in conversational chain: {str(e)}"
        print(error_msg)
        return f"I apologize, but I encountered an error while processing your question: {str(e)}"

def process_user_question(user_question, chat_history):
    """Process user question and return response with retrieval results"""
    try:
        new_db = FAISS.load_local("faiss_db", embeddings, allow_dangerous_deserialization=True)
        # retriever = new_db.as_retriever()
        
        # Get top 5 retrieval results for this specific question
        top_results = get_top_retrieval_results(new_db, user_question, k=5)
        
        # Get response using direct RAG approach
        response = get_conversational_chain(new_db, user_question, chat_history)
        
        return response, top_results
        
    except Exception as e:
        return f"Error: Unable to process your question. Please make sure you have uploaded and processed files first. Details: {str(e)}", []

def display_retrieval_results(results, question):
    """Display retrieval results for a specific question"""
    if results:
        print(results)
        with st.expander(f"📄 Top 5 Retrieved Sections for: '{question[:50]}...'", expanded=False):
            for i, doc in enumerate(results, 1):
                st.write(f"**Result {i}:**")
                
                # Handle both dictionary and object formats
                if isinstance(doc, dict):
                    # For dictionary format (your current JSON structure)
                    content = doc.get('content', '')
                    source = doc.get('source', 'Unknown source')
                    score = doc.get('score', 0)
                    chunk_index = doc.get('chunk_index', 'N/A')
                    total_chunks = doc.get('total_chunks', 'N/A')
                    
                    # Truncate content if too long
                    display_content = content[:300] + "..." if len(content) > 300 else content
                    
                    st.write(display_content)
                    st.write(f"**Source:** {source}")
                    st.write(f"**Score:** {score:.4f}")
                    st.write(f"**Chunk:** {chunk_index}/{total_chunks}")
                    
                else:
                    # For object format (original structure with page_content attribute)
                    content = doc.page_content[:300] + "..." if len(doc.page_content) > 300 else doc.page_content
                    st.write(content)
                
                if i < len(results):
                    st.write("---")

def get_file_icon(filename):
    """Return appropriate icon for file type"""
    filename = filename.lower()
    if filename.endswith('.pdf'):
        return "📄"
    elif filename.endswith('.docx'):
        return "📝"
    elif filename.endswith('.txt'):
        return "📄"
    elif filename.endswith('.csv'):
        return "📊"
    elif filename.endswith(('.xlsx', '.xls')):
        return "📈"
    elif filename.endswith(('.pptx', '.ppt')):
        return "📽️"
    else:
        return "📁"

def main():
    st.set_page_config(
        page_title="ARD File Chatbot",
        page_icon="🤖",
        layout="wide"
    )
    
    # Initialize session state
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    
    if 'files_processed' not in st.session_state:
        st.session_state.files_processed = False
    
    if 'processed_files_info' not in st.session_state:
        st.session_state.processed_files_info = []

    st.title("🤖 ARD File Chatbot")
    st.markdown("Upload your documents and start chatting! Supports: PDF, DOCX, TXT, CSV, Excel, PowerPoint, R, Python, Java, C/C++")

    # Main layout
    col1, col2 = st.columns([2, 1])
    
    with col1:
        # Chat interface
        if not st.session_state.files_processed:
            st.warning("⚠️ Please upload and process your files first before asking questions!")
        else:
            st.success(f"✅ Files processed successfully! You can now ask questions about your {len(st.session_state.processed_files_info)} uploaded files.")
        
        # Chat input at the top for better UX
        user_question = st.chat_input(
            "Ask a question about your documents...",
            disabled=not st.session_state.files_processed
        )
        
        # Process new question
        if user_question and st.session_state.files_processed:
            # Add user message to chat history
            st.session_state.chat_history.append({
                "role": "user", 
                "content": user_question
            })
            
            # Get response and retrieval results
            with st.spinner("Thinking..."):
                response, retrieval_results = process_user_question(user_question, st.session_state.chat_history[:-1])
            
            if "answer is not available in the context" in response.lower() or "not available in the context" in response.lower():
                retrieval_results = []

            # Add assistant response to chat history
            st.session_state.chat_history.append({
                "role": "assistant", 
                "content": response,
                "retrieval_results": retrieval_results,
                "question": user_question
            })
        
        # Display chat history with retrieval results
        st.subheader("💬 Conversation")
        
        if st.session_state.chat_history:
            # Create a scrollable container for chat
            chat_container = st.container()
            
            with chat_container:
                for i, message in enumerate(st.session_state.chat_history):
                    if message["role"] == "user":
                        with st.chat_message("user"):
                            st.write(message["content"])
                    else:
                        with st.chat_message("assistant"):
                            st.write(message["content"])
                            
                            # Show retrieval results only for this specific question-answer pair
                            if "retrieval_results" in message and message["retrieval_results"]:
                                display_retrieval_results(
                                    message["retrieval_results"], 
                                    message.get("question", "Question")
                                )
        else:
            if st.session_state.files_processed:
                st.info("👋 Start by asking a question about your uploaded documents!")
            else:
                st.info("📁 Upload and process your files to begin chatting!")
    
    with col2:
        # File management sidebar
        st.subheader("📁 File Management")
        
        # Display supported formats
        st.markdown("**Supported Formats:**")
        st.markdown("📄 PDF • 📝 DOCX • 📄 TXT • 📊 CSV • 📈 Excel • 📽️ PowerPoint • 🐍 Python • 📊 R • ⚡ C/C++ • ☕ Java")
        st.write("---")
        
        # Show processed files info
        if st.session_state.processed_files_info:
            st.write("**Processed Files:**")
            for file_info in st.session_state.processed_files_info:
                st.write(f"• {file_info}")
            st.write("---")
        
        # File upload
        files = st.file_uploader(
            "Upload your files here and click on 'Process'", 
            accept_multiple_files=True,
            type=SUPPORTED_EXTENSIONS,
            help="Supported formats: PDF, DOCX, TXT, CSV, XLSX, XLS, PPTX, PPT, R, PYTHON, C/C++, JAVA"
        )
        
        # Process files button
        if st.button("🔄 Process Files", type="primary", use_container_width=True):
            if not files:
                st.error("Please upload at least one file to process.")
            else:
                with st.spinner("Processing files..."):
                    try:
                        # Process different file types
                        all_chunks_with_metadata = []
                        processed_files = []
                        error_files = []
                        document_registry = {}  # Keep track of processed documents
                        
                        for file_index, file in enumerate(files):
                            # Check if file is supported
                            if not is_supported_file(file.name):
                                error_files.append(f"❌ {file.name} (unsupported format)")
                                continue
                            
                            # Create unique document ID
                            document_id = f"doc_{file_index}_{hash(file.name) % 10000}"
                            
                            # Extract text based on file type
                            text = get_file_text(file)
                            
                            if text and not text.startswith("Error"):
                                # Create chunks with metadata for this document
                                chunks_with_metadata = get_chunks_with_metadata(text, document_id, file.name)
                                all_chunks_with_metadata.extend(chunks_with_metadata)
                                
                                # Store document info
                                document_registry[document_id] = {
                                    'name': file.name,
                                    'chunks_count': len(chunks_with_metadata),
                                    'file_type': file.name.split('.')[-1].upper()
                                }
                                
                                icon = get_file_icon(file.name)
                                processed_files.append(f"{icon} {file.name} ({len(chunks_with_metadata)} chunks)")
                            else:
                                error_files.append(f"❌ {file.name} (processing error)")
                                
                        # Show errors if any
                        if error_files:
                            st.error("Some files could not be processed:")
                            for error_file in error_files:
                                st.write(f"  {error_file}")
                        
                        if all_chunks_with_metadata:
                            vector_store = create_vector_store_with_metadata(all_chunks_with_metadata)
                            
                            # Update session state
                            st.session_state.files_processed = True
                            st.session_state.processed_files_info = processed_files
                            st.session_state.document_registry = document_registry
                            st.session_state.total_chunks = len(all_chunks_with_metadata)
                            st.session_state.chat_history = []  # Reset chat for new files
                            
                            # Display processing summary
                            st.success(f"✅ Successfully processed {len(processed_files)} files!")
                            
                            # Show document summary
                            with st.expander("📊 Processing Summary", expanded=True):
                                col1, col2 = st.columns(2)
                                with col1:
                                    st.metric("Total Documents", len(document_registry))
                                    st.metric("Total Chunks", len(all_chunks_with_metadata))
                                with col2:
                                    for doc_id, info in document_registry.items():
                                        st.write(f"📄 **{info['name']}**: {info['chunks_count']} chunks")
                            
                            st.rerun()
                        else:
                            st.error("No files were successfully processed. Please check your file formats and try again.")
                            
                    except Exception as e:
                        st.error(f"Error processing files: {str(e)}")
        
        st.write("---")
        
        # Action buttons
        col_a, col_b = st.columns(2)
        
        with col_a:
            if st.button("🗑️ Clear Chat", use_container_width=True):
                if st.session_state.chat_history:
                    st.session_state.chat_history = []
                    st.rerun()
        
        with col_b:
            if st.button("🔄 Reset All", use_container_width=True):
                if st.session_state.files_processed:
                    st.session_state.files_processed = False
                    st.session_state.processed_files_info = []
                    st.session_state.chat_history = []
                    # Clean up vector store
                    try:
                        import shutil
                        if os.path.exists("faiss_db"):
                            shutil.rmtree("faiss_db")
                    except:
                        pass
                    st.rerun()
        
        # Chat statistics
        if st.session_state.chat_history:
            st.write("---")
            st.write("**Chat Statistics:**")
            user_messages = len([msg for msg in st.session_state.chat_history if msg["role"] == "user"])
            st.write(f"• Questions asked: {user_messages}")
            st.write(f"• Total messages: {len(st.session_state.chat_history)}")

if __name__ == "__main__":
    main()